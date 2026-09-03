from pptx import Presentation

prs = Presentation(r"c:\src\projects\SIH\voiceguard\reference\SIH2025-IDEA-Presentation-Format.pptx")

for i, slide in enumerate(prs.slides):
    print(f"\n=================== SLIDE {i+1} ===================")
    for shape in slide.shapes:
        if shape.has_text_frame:
            print(f"Shape: {shape.name} | type={shape.shape_type}")
            for p in shape.text_frame.paragraphs:
                for r in p.runs:
                    color_str = "None"
                    if r.font and r.font.color:
                        try:
                            color_str = f"rgb:{r.font.color.rgb}"
                        except Exception:
                            try:
                                color_str = f"theme:{r.font.color.theme_color}"
                            except Exception:
                                color_str = "other"
                    print(f"   Run: '{r.text.strip()}' | font={r.font.name}, size={r.font.size.pt if r.font.size else None}, bold={r.font.bold}, color={color_str}")
