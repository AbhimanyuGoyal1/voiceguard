from pptx import Presentation

prs = Presentation(r"c:\src\projects\SIH\voiceguard\reference\SIH2025-IDEA-Presentation-Format.pptx")

for idx, slide in enumerate(prs.slides):
    print(f"=== SLIDE {idx+1} ===")
    for shape in slide.shapes:
        left = shape.left.inches if shape.left else 0
        top = shape.top.inches if shape.top else 0
        width = shape.width.inches if shape.width else 0
        height = shape.height.inches if shape.height else 0
        print(f"Shape '{shape.name}' ({shape.shape_type}) pos: ({left:.2f}, {top:.2f}, {width:.2f}, {height:.2f})")
        if shape.has_text_frame:
            for p in shape.text_frame.paragraphs:
                font_name = p.font.name if p.font else None
                font_size = p.font.size.pt if p.font and p.font.size else None
                color = None
                if p.font and p.font.color and hasattr(p.font.color, 'rgb'):
                    color = p.font.color.rgb
                print(f"   P: font={font_name}, size={font_size}, color={color} -> '{p.text}'")
