import sys
from pptx import Presentation
from pptx.enum.shapes import MSO_SHAPE_TYPE

prs = Presentation(r"c:\src\projects\SIH\voiceguard\reference\SIH2025-IDEA-Presentation-Format.pptx")

print(f"Slide width: {prs.slide_width.inches} in, Slide height: {prs.slide_height.inches} in")
print(f"Number of slides: {len(prs.slides)}")

for idx, slide in enumerate(prs.slides):
    print(f"\n--- SLIDE {idx+1} ---")
    layout = slide.slide_layout
    print(f"Layout name: {layout.name if layout else 'None'}")
    for shape in slide.shapes:
        info = f"Shape: name='{shape.name}', type={shape.shape_type}"
        if shape.has_text_frame:
            text = shape.text_frame.text.replace('\n', ' \\n ')
            print(f"  {info} | Text: {text[:100]}")
        else:
            print(f"  {info}")
