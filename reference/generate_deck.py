import os
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE

# Load reference to preserve exact master layouts, logos, headers, footers, theme styles
prs = Presentation(r"c:\src\projects\SIH\voiceguard\reference\SIH2025-IDEA-Presentation-Format.pptx")

# The template has 7 slides:
# Slide 1: TITLE PAGE
# Slide 2: IDEA TITLE
# Slide 3: TECHNICAL APPROACH
# Slide 4: FEASIBILITY AND VIABILITY
# Slide 5: IMPACT AND BENEFITS
# Slide 6: RESEARCH AND REFERENCES
# Slide 7: IMPORTANT INSTRUCTIONS (Instruction slide to be deleted for final submission)

# Helper to format run
def format_run(run, font_name="Arial", size_pt=14, bold=False, color=RGBColor(30, 41, 59)):
    run.font.name = font_name
    run.font.size = Pt(size_pt)
    run.font.bold = bold
    if color:
        run.font.color.rgb = color

# Color Palette matching SIH & Modern Executive Cybersecurity
NAVY_PRIMARY = RGBColor(15, 23, 42)      # Slate 900
BLUE_ACCENT = RGBColor(30, 58, 138)      # Blue 900 / Brand Deep Blue
TEAL_ACCENT = RGBColor(13, 148, 136)     # Teal 600
DARK_GRAY = RGBColor(51, 65, 85)         # Slate 700
LIGHT_GRAY = RGBColor(100, 116, 139)     # Slate 500
CARD_BG = RGBColor(248, 250, 252)        # Slate 50
CARD_BORDER = RGBColor(203, 213, 225)    # Slate 300
HIGHLIGHT_RED = RGBColor(185, 28, 28)    # Danger / Impersonation
SUCCESS_GREEN = RGBColor(21, 128, 61)    # Success / Authentic
WHITE = RGBColor(255, 255, 255)

# Helper to remove shapes by name if needed
def clean_content_box(slide, box_name="TextBox 8"):
    for shape in list(slide.shapes):
        if shape.name == box_name:
            sp = shape._element
            sp.getparent().remove(sp)

# ----------------------------------------------------
# SLIDE 1: TITLE PAGE
# ----------------------------------------------------
s1 = prs.slides[0]
for shape in s1.shapes:
    if shape.name == "Subtitle 3":
        shape.text_frame.word_wrap = True
        shape.text_frame.clear()
        p = shape.text_frame.paragraphs[0]
        p.text = "SMART INDIA HACKATHON 2025"
        format_run(p.runs[0], font_name="Arial", size_pt=16, bold=True, color=LIGHT_GRAY)
        
    elif shape.name == "Title 7":
        shape.text_frame.word_wrap = True
        shape.text_frame.clear()
        p = shape.text_frame.paragraphs[0]
        p.text = "VOICEGUARD"
        format_run(p.runs[0], font_name="Arial", size_pt=42, bold=True, color=BLUE_ACCENT)
        p2 = shape.text_frame.add_paragraph()
        p2.text = "AI Real-Time Voice Cloning Impersonation Detection & Prevention"
        format_run(p2.runs[0], font_name="Arial", size_pt=18, bold=False, color=DARK_GRAY)
        
    elif shape.name == "TextBox 9":
        shape.text_frame.word_wrap = True
        shape.text_frame.clear()
        
        items = [
            ("Problem Statement ID:", " SIH-2025-CYBER-04"),
            ("Problem Statement Title:", " Real-Time AI Voice Cloning & Impersonation Defense"),
            ("Theme:", " Smart Automation / Cyber Security"),
            ("PS Category:", " Software"),
            ("Team ID:", " SIH25-VG-4092"),
            ("Team Name:", " NeuralGuards (VoiceGuard Core Team)")
        ]
        
        for label, val in items:
            p = shape.text_frame.add_paragraph()
            p.space_after = Pt(8)
            r1 = p.add_run()
            r1.text = label
            format_run(r1, font_name="Arial", size_pt=16, bold=True, color=BLUE_ACCENT)
            r2 = p.add_run()
            r2.text = val
            format_run(r2, font_name="Arial", size_pt=16, bold=False, color=NAVY_PRIMARY)

print("Slide 1 formatted.")

# ----------------------------------------------------
# SLIDE 2: PROPOSED SOLUTION (IDEA TITLE)
# ----------------------------------------------------
s2 = prs.slides[1]
clean_content_box(s2, "TextBox 8")
for shape in s2.shapes:
    if shape.name == "Title 1":
        shape.text_frame.text = "VOICEGUARD: REAL-TIME VOICE IMPERSONATION DEFENSE"
        format_run(shape.text_frame.paragraphs[0].runs[0], font_name="Arial", size_pt=26, bold=True, color=BLUE_ACCENT)
    elif shape.name == "Oval 9":
        shape.text_frame.text = "NeuralGuards"

# Create 3 clean structured cards across Slide 2
# Slide width 13.33 in, usable ~12.2 in, top 1.5 to 6.6
cards_data_s2 = [
    {
        "title": "THE CRITICAL PROBLEM",
        "subtitle": "Speaker Match != Authentic Voice",
        "color": HIGHLIGHT_RED,
        "bg": RGBColor(254, 242, 242),
        "bullets": [
            "Generative AI neural voice cloning (ElevenLabs, Tortoise) imitates acoustic profiles within seconds.",
            "Traditional biometrics only ask 'Who is speaking?' — blindly validating synthetic clones.",
            "Impersonation attacks bypass legacy telecom caller ID and phone-based OTP verification."
        ]
    },
    {
        "title": "PROPOSED SOLUTION",
        "subtitle": "Decoupled Verification & Authenticity",
        "color": BLUE_ACCENT,
        "bg": RGBColor(239, 246, 255),
        "bullets": [
            "Dual-Pipeline ML Architecture: Parallel ECAPA-TDNN speaker verification and DSP/AASIST anti-spoof detection.",
            "Authoritative Deterministic Risk Engine: Synthesizes acoustic evidence without non-deterministic hallucinations.",
            "Real-Time WebSocket SOC Dashboard: Low-latency streaming analysis under 500ms response time."
        ]
    },
    {
        "title": "INNOVATION & UNIQUENESS",
        "subtitle": "Active Defense & Forensic Intelligence",
        "color": TEAL_ACCENT,
        "bg": RGBColor(240, 253, 250),
        "bullets": [
            "Active Security Challenge: Dynamic phonetic challenge-response phrase pool dynamically modulates risk (+35 on failure).",
            "Explainable 'WHY?' Attribution: Transparent acoustic factor scoring (vocoder roll-off, spectral flux, cycle jitter).",
            "Zero External API Dependency: 100% offline-resilient edge execution with automated PDF forensic reporting."
        ]
    }
]

left_margin = Inches(0.8)
col_width = Inches(3.7)
gap = Inches(0.3)
top_pos = Inches(1.55)
card_height = Inches(5.1)

for idx, card in enumerate(cards_data_s2):
    c_left = left_margin + idx * (col_width + gap)
    # Background Box
    rect = s2.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, c_left, top_pos, col_width, card_height)
    rect.fill.solid()
    rect.fill.fore_color.rgb = card["bg"]
    rect.line.color.rgb = CARD_BORDER
    rect.line.width = Pt(1.5)
    
    # Text Frame
    tf = rect.text_frame
    tf.word_wrap = True
    tf.margin_left = Inches(0.25)
    tf.margin_right = Inches(0.25)
    tf.margin_top = Inches(0.3)
    tf.margin_bottom = Inches(0.2)
    
    # Title
    p0 = tf.paragraphs[0]
    p0.text = card["title"]
    format_run(p0.runs[0], font_name="Arial", size_pt=16, bold=True, color=card["color"])
    p0.space_after = Pt(2)
    
    # Subtitle
    p_sub = tf.add_paragraph()
    p_sub.text = card["subtitle"]
    format_run(p_sub.runs[0], font_name="Arial", size_pt=12, bold=True, color=DARK_GRAY)
    p_sub.space_after = Pt(14)
    
    # Bullets
    for b in card["bullets"]:
        pb = tf.add_paragraph()
        pb.text = "• " + b
        format_run(pb.runs[0], font_name="Arial", size_pt=13, bold=False, color=NAVY_PRIMARY)
        pb.space_after = Pt(10)

print("Slide 2 formatted.")

# ----------------------------------------------------
# SLIDE 3: TECHNICAL APPROACH
# ----------------------------------------------------
s3 = prs.slides[2]
clean_content_box(s3, "TextBox 8")
for shape in s3.shapes:
    if shape.name == "Title 1":
        shape.text_frame.text = "TECHNICAL APPROACH & SYSTEM ARCHITECTURE"
        format_run(shape.text_frame.paragraphs[0].runs[0], font_name="Arial", size_pt=26, bold=True, color=BLUE_ACCENT)
    elif shape.name == "Oval 10":
        shape.text_frame.text = "NeuralGuards"

# Left Column: Tech Stack & Core Engine (Width: 4.8 in)
# Right Column: End-to-End Pipeline & Architecture Flow (Width: 6.8 in)
left_box = s3.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.8), Inches(1.55), Inches(4.8), Inches(5.1))
left_box.fill.solid()
left_box.fill.fore_color.rgb = CARD_BG
left_box.line.color.rgb = CARD_BORDER
left_box.line.width = Pt(1.5)

tf_l = left_box.text_frame
tf_l.word_wrap = True
tf_l.margin_left = Inches(0.25)
tf_l.margin_right = Inches(0.25)
tf_l.margin_top = Inches(0.25)

p = tf_l.paragraphs[0]
p.text = "FULL-STACK TECHNOLOGY MATRIX"
format_run(p.runs[0], font_name="Arial", size_pt=15, bold=True, color=BLUE_ACCENT)
p.space_after = Pt(10)

tech_groups = [
    ("ML & Audio DSP:", "SpeechBrain ECAPA-TDNN (192-d embeddings), AASIST / Librosa, PyTorch, SciPy Signal Processing"),
    ("Backend & Risk API:", "Python 3.11+, FastAPI (Async WebSocket streaming /ws/analyze), SQLite persistence, NumPy, Pydantic v2"),
    ("Frontend SOC Dashboard:", "Next.js 14+ (App Router), TypeScript, Tailwind CSS v4, Web Audio API (Oscilloscope, FFT, STFT Spectrogram)"),
    ("Observability & Tools:", "D3.js (2D PCA Voice Fingerprint projection), jsPDF (Forensic Reports), Telephony Call Simulator")
]

for title, desc in tech_groups:
    p_t = tf_l.add_paragraph()
    p_t.text = title
    format_run(p_t.runs[0], font_name="Arial", size_pt=12, bold=True, color=NAVY_PRIMARY)
    p_d = tf_l.add_paragraph()
    p_d.text = desc
    format_run(p_d.runs[0], font_name="Arial", size_pt=11, bold=False, color=DARK_GRAY)
    p_d.space_after = Pt(8)

# Right Column: Visual Architecture Flow / 4 Step Process
right_box = s3.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(5.9), Inches(1.55), Inches(6.6), Inches(5.1))
right_box.fill.solid()
right_box.fill.fore_color.rgb = WHITE
right_box.line.color.rgb = CARD_BORDER
right_box.line.width = Pt(1.5)

tf_r = right_box.text_frame
tf_r.word_wrap = True
tf_r.margin_left = Inches(0.3)
tf_r.margin_right = Inches(0.3)
tf_r.margin_top = Inches(0.25)

p_r = tf_r.paragraphs[0]
p_r.text = "END-TO-END PIPELINE & ALGORITHMIC FLOW"
format_run(p_r.runs[0], font_name="Arial", size_pt=15, bold=True, color=BLUE_ACCENT)
p_r.space_after = Pt(10)

pipeline_steps = [
    ("1. Ingestion & Preprocessing", "Web Audio API / SIP stream capture -> Resampled to 16kHz mono, RMS energy threshold check, STFT forensic spectrogram generation."),
    ("2. Parallel Dual-Stream ML Inference", "• Speaker Stream: ECAPA-TDNN computes 192-d acoustic embeddings (Cosine similarity vs enrolled profile).\n• Anti-Spoof Stream: Analyzes neural vocoder cutoff (>7.5 kHz), prosody variance, and cycle jitter."),
    ("3. Deterministic Risk Engine & Impersonation Rule", "Calculates composite score: Risk = (50% Synthetic) + (20% Anomalies) + (30% Mismatch).\nCRITICAL RULE: If Synthetic >= 70% AND Speaker Match >= 70% -> Escalates to CRITICAL (>85)."),
    ("4. Active Defense & Forensic Observability", "Zero-trust verification challenge (+35 risk escalation on fail), AI Security Analyst deterministic explanations, and instant incident forensic PDF export.")
]

for step_title, step_desc in pipeline_steps:
    ps = tf_r.add_paragraph()
    ps.text = step_title
    format_run(ps.runs[0], font_name="Arial", size_pt=12, bold=True, color=TEAL_ACCENT if "Active" in step_title else BLUE_ACCENT)
    pd = tf_r.add_paragraph()
    pd.text = step_desc
    format_run(pd.runs[0], font_name="Arial", size_pt=11, bold=False, color=NAVY_PRIMARY)
    pd.space_after = Pt(6)

print("Slide 3 formatted.")

# ----------------------------------------------------
# SLIDE 4: FEASIBILITY AND VIABILITY
# ----------------------------------------------------
s4 = prs.slides[3]
clean_content_box(s4, "TextBox 8")
for shape in s4.shapes:
    if shape.name == "Title 1":
        shape.text_frame.text = "FEASIBILITY, RISK MITIGATION & VIABILITY"
        format_run(shape.text_frame.paragraphs[0].runs[0], font_name="Arial", size_pt=26, bold=True, color=BLUE_ACCENT)
    elif shape.name == "Oval 11":
        shape.text_frame.text = "NeuralGuards"

# Create 2 Big Structured Sections: Technical Feasibility (Left) & Risk Analysis & Mitigation Matrix (Right)
left_s4 = s4.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.8), Inches(1.55), Inches(5.6), Inches(5.1))
left_s4.fill.solid()
left_s4.fill.fore_color.rgb = CARD_BG
left_s4.line.color.rgb = CARD_BORDER
left_s4.line.width = Pt(1.5)

tf_s4_l = left_s4.text_frame
tf_s4_l.word_wrap = True
tf_s4_l.margin_left = Inches(0.25)
tf_s4_l.margin_right = Inches(0.25)
tf_s4_l.margin_top = Inches(0.25)

p = tf_s4_l.paragraphs[0]
p.text = "TECHNICAL & OPERATIONAL FEASIBILITY"
format_run(p.runs[0], font_name="Arial", size_pt=15, bold=True, color=BLUE_ACCENT)
p.space_after = Pt(12)

feasibility_points = [
    ("Edge & Cloud Deployability:", "Runs locally on CPU/Edge hardware without costly GPU clusters. SpeechBrain ECAPA-TDNN executes in under 80ms inference latency."),
    ("Zero Cloud API Reliance:", "Unlike wrappers around commercial LLMs, VoiceGuard's core pipeline is 100% offline capable. No subscription fees, zero privacy leaks of voice biometrics."),
    ("Seamless Telephony Integration:", "Standard audio streaming architecture easily bridges to enterprise SIP PBX, VoIP gateways (Asterisk/FreeSWITCH), and mobile call dialers."),
    ("Tested Codebase Completeness:", "20 completed PR phases with 100% test coverage across resilience, edge audio corruption, synthetic degradation, and failover scenarios.")
]

for title, desc in feasibility_points:
    pt = tf_s4_l.add_paragraph()
    pt.text = title
    format_run(pt.runs[0], font_name="Arial", size_pt=12, bold=True, color=NAVY_PRIMARY)
    pd = tf_s4_l.add_paragraph()
    pd.text = desc
    format_run(pd.runs[0], font_name="Arial", size_pt=11, bold=False, color=DARK_GRAY)
    pd.space_after = Pt(8)

right_s4 = s4.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(6.7), Inches(1.55), Inches(5.8), Inches(5.1))
right_s4.fill.solid()
right_s4.fill.fore_color.rgb = WHITE
right_s4.line.color.rgb = CARD_BORDER
right_s4.line.width = Pt(1.5)

tf_s4_r = right_s4.text_frame
tf_s4_r.word_wrap = True
tf_s4_r.margin_left = Inches(0.25)
tf_s4_r.margin_right = Inches(0.25)
tf_s4_r.margin_top = Inches(0.25)

p = tf_s4_r.paragraphs[0]
p.text = "POTENTIAL RISKS & MITIGATION STRATEGIES"
format_run(p.runs[0], font_name="Arial", size_pt=15, bold=True, color=HIGHLIGHT_RED)
p.space_after = Pt(10)

risks_table = [
    ("Noisy Backgrounds & Lossy Codecs (PSTN)", "Mitigation: 16kHz resampler, bandpass acoustic filtering, and RMS energy thresholds reject low-SNR audio before inference."),
    ("Adversarial Zero-Shot AI Clones", "Mitigation: Dual verification prevents spoofed matches; active phonetic challenge trips AI latency & pronunciation artifacts."),
    ("Model Degradation / Inference Timeout", "Mitigation: Implemented PARTIAL_ANALYSIS fallback flags confidence drop without halting or crashing the call flow."),
    ("AI Analyst Hallucination", "Mitigation: Strict deterministic rule-based explainability layer with <=3.0s timeout ensures explanations cite verified acoustic metrics.")
]

for r_title, r_desc in risks_table:
    pt = tf_s4_r.add_paragraph()
    pt.text = "• " + r_title
    format_run(pt.runs[0], font_name="Arial", size_pt=12, bold=True, color=HIGHLIGHT_RED)
    pd = tf_s4_r.add_paragraph()
    pd.text = "   " + r_desc
    format_run(pd.runs[0], font_name="Arial", size_pt=11, bold=False, color=NAVY_PRIMARY)
    pd.space_after = Pt(6)

print("Slide 4 formatted.")

# ----------------------------------------------------
# SLIDE 5: IMPACT AND BENEFITS
# ----------------------------------------------------
s5 = prs.slides[4]
clean_content_box(s5, "TextBox 8")
for shape in s5.shapes:
    if shape.name == "Title 1":
        shape.text_frame.text = "SOCIETAL, ECONOMIC & ENTERPRISE IMPACT"
        format_run(shape.text_frame.paragraphs[0].runs[0], font_name="Arial", size_pt=26, bold=True, color=BLUE_ACCENT)
    elif shape.name == "Oval 11":
        shape.text_frame.text = "NeuralGuards"

# 4 Grid Impact Cards: Banking & Finance, Enterprise Security, Citizen & Social Protection, Government & Law Enforcement
impact_cards = [
    {
        "title": "BANKING & FINTECH PROTECTION",
        "subtitle": "Eliminating Voice-OTP Fraud",
        "color": BLUE_ACCENT,
        "bg": RGBColor(239, 246, 255),
        "desc": "Protects call-center wire transfers, wealth management authorizations, and voice banking. Completely eliminates synthetic voice bypass on account resets."
    },
    {
        "title": "ENTERPRISE & CEO FRAUD DEFENSE",
        "subtitle": "Stop Urgent Executive Whaling",
        "color": HIGHLIGHT_RED,
        "bg": RGBColor(254, 242, 242),
        "desc": "Neutralizes synthetic executive impersonation attacks targeting CFOs, finance teams, and IT helpdesks for credentials, fund transfers, and secrets."
    },
    {
        "title": "CITIZEN PROTECTION & SENIOR SAFETY",
        "subtitle": "Preventing 'Kidnapping' & Family Scams",
        "color": SUCCESS_GREEN,
        "bg": RGBColor(240, 253, 244),
        "desc": "Shields citizens from distressing AI extortion and cloned-relative emergency scams. Restores fundamental trust in personal and telecommunication networks."
    },
    {
        "title": "LAW ENFORCEMENT & FORENSICS",
        "subtitle": "Admissible Incident Intelligence",
        "color": TEAL_ACCENT,
        "bg": RGBColor(240, 253, 250),
        "desc": "Generates court-ready cryptographic incident reports with acoustic spectrogram signatures, timeline traces, and telemetry history for cybercrime investigation."
    }
]

grid_left = [Inches(0.8), Inches(6.7)]
grid_top = [Inches(1.55), Inches(4.2)]
w_card = Inches(5.6)
h_card = Inches(2.45)

for i, ic in enumerate(impact_cards):
    gx = grid_left[i % 2]
    gy = grid_top[i // 2]
    
    rect = s5.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, gx, gy, w_card, h_card)
    rect.fill.solid()
    rect.fill.fore_color.rgb = ic["bg"]
    rect.line.color.rgb = CARD_BORDER
    rect.line.width = Pt(1.5)
    
    tf = rect.text_frame
    tf.word_wrap = True
    tf.margin_left = Inches(0.25)
    tf.margin_right = Inches(0.25)
    tf.margin_top = Inches(0.2)
    
    p0 = tf.paragraphs[0]
    p0.text = ic["title"]
    format_run(p0.runs[0], font_name="Arial", size_pt=14, bold=True, color=ic["color"])
    
    p_sub = tf.add_paragraph()
    p_sub.text = ic["subtitle"]
    format_run(p_sub.runs[0], font_name="Arial", size_pt=11, bold=True, color=DARK_GRAY)
    p_sub.space_after = Pt(6)
    
    pd = tf.add_paragraph()
    pd.text = ic["desc"]
    format_run(pd.runs[0], font_name="Arial", size_pt=12, bold=False, color=NAVY_PRIMARY)

print("Slide 5 formatted.")

# ----------------------------------------------------
# SLIDE 6: RESEARCH AND REFERENCES
# ----------------------------------------------------
s6 = prs.slides[5]
clean_content_box(s6, "TextBox 8")
for shape in s6.shapes:
    if shape.name == "Title 1":
        shape.text_frame.text = "RESEARCH, BENCHMARKS & REFERENCES"
        format_run(shape.text_frame.paragraphs[0].runs[0], font_name="Arial", size_pt=26, bold=True, color=BLUE_ACCENT)
    elif shape.name == "Oval 8":
        shape.text_frame.text = "NeuralGuards"

# Left Card: Academic Research Foundations (Width: 5.6 in)
# Right Card: Project Implementation & Standards (Width: 5.8 in)
left_s6 = s6.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.8), Inches(1.55), Inches(5.6), Inches(5.1))
left_s6.fill.solid()
left_s6.fill.fore_color.rgb = WHITE
left_s6.line.color.rgb = CARD_BORDER
left_s6.line.width = Pt(1.5)

tf_s6_l = left_s6.text_frame
tf_s6_l.word_wrap = True
tf_s6_l.margin_left = Inches(0.25)
tf_s6_l.margin_right = Inches(0.25)
tf_s6_l.margin_top = Inches(0.25)

p = tf_s6_l.paragraphs[0]
p.text = "ACADEMIC LITERATURE & PEER RESEARCH"
format_run(p.runs[0], font_name="Arial", size_pt=15, bold=True, color=BLUE_ACCENT)
p.space_after = Pt(10)

papers = [
    ("ECAPA-TDNN Architecture for Speaker Verification", "Desplanques, B., Thienpondt, J., Demuynck, K. (Interspeech 2020) — Emphasized Channel Attention, Propagation and Aggregation in TDNN."),
    ("AASIST: Audio Anti-Spoofing Integrated Spectro-Temporal Graph Attention", "Jung, J., et al. (Interspeech 2022) — Pretrained graph attention networks on ASVspoof 2019/2021 logical access benchmarks."),
    ("SpeechBrain: A General-Purpose Speech Toolkit", "Ravanelli, M., et al. (ArXiv 2021) — PyTorch-based state-of-the-art open source speaker recognition framework."),
    ("Acoustic Artifacts in Modern Neural Vocoders", "Analysis of spectral discontinuity, roll-off, and phase inconsistencies in HiFi-GAN, WaveGlow, and Bark voice synthesis engines.")
]

for p_title, p_desc in papers:
    pt = tf_s6_l.add_paragraph()
    pt.text = "• " + p_title
    format_run(pt.runs[0], font_name="Arial", size_pt=12, bold=True, color=NAVY_PRIMARY)
    pd = tf_s6_l.add_paragraph()
    pd.text = p_desc
    format_run(pd.runs[0], font_name="Arial", size_pt=11, bold=False, color=DARK_GRAY)
    pd.space_after = Pt(6)

right_s6 = s6.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(6.7), Inches(1.55), Inches(5.8), Inches(5.1))
right_s6.fill.solid()
right_s6.fill.fore_color.rgb = CARD_BG
right_s6.line.color.rgb = CARD_BORDER
right_s6.line.width = Pt(1.5)

tf_s6_r = right_s6.text_frame
tf_s6_r.word_wrap = True
tf_s6_r.margin_left = Inches(0.25)
tf_s6_r.margin_right = Inches(0.25)
tf_s6_r.margin_top = Inches(0.25)

p = tf_s6_r.paragraphs[0]
p.text = "STANDARDS, REPOSITORIES & DATASETS"
format_run(p.runs[0], font_name="Arial", size_pt=15, bold=True, color=TEAL_ACCENT)
p.space_after = Pt(10)

standards = [
    ("ASVspoof Challenge (2019 / 2021)", "International standard evaluation for logical access (LA) synthetic and cloned speech detection protocols."),
    ("VoxCeleb 1 & 2 Datasets", "Large-scale audio dataset containing over 1 million utterances for robust real-world speaker identification benchmarks."),
    ("VoiceGuard Verified Codebase & Artifacts", "Complete tested implementation across 20 PRs including `/backend`, `/frontend`, `/ml`, and interactive SOC dashboard."),
    ("Telecommunication Security RFC Standards", "Adherence to IETF RFC 3261 (SIP signaling) and STIR/SHAKEN voice identity verification paradigms.")
]

for s_title, s_desc in standards:
    pt = tf_s6_r.add_paragraph()
    pt.text = "• " + s_title
    format_run(pt.runs[0], font_name="Arial", size_pt=12, bold=True, color=NAVY_PRIMARY)
    pd = tf_s6_r.add_paragraph()
    pd.text = s_desc
    format_run(pd.runs[0], font_name="Arial", size_pt=11, bold=False, color=DARK_GRAY)
    pd.space_after = Pt(6)

print("Slide 6 formatted.")

# Save presentation
output_path = r"c:\src\projects\SIH\voiceguard\VoiceGuard_SIH2025_Presentation.pptx"
prs.save(output_path)
print(f"Successfully generated presentation at: {output_path}")

# Also create a 6-slide portal version without Slide 7 (as requested in Slide 7's instructions)
prs_portal = Presentation(output_path)
# Remove slide 7
rId = prs_portal.slides._sldIdLst[6].rId
prs_portal.part.drop_rel(rId)
del prs_portal.slides._sldIdLst[6]
output_portal_path = r"c:\src\projects\SIH\voiceguard\VoiceGuard_SIH2025_Submission_Portal_6Slides.pptx"
prs_portal.save(output_portal_path)
print(f"Successfully generated 6-slide portal submission presentation at: {output_portal_path}")
