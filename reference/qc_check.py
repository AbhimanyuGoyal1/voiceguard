from pptx import Presentation

for path in [
    r"c:\src\projects\SIH\voiceguard\VoiceGuard_SIH2025_Presentation.pptx",
    r"c:\src\projects\SIH\voiceguard\VoiceGuard_SIH2025_Submission_Portal_6Slides.pptx"
]:
    prs = Presentation(path)
    print(f"\n==========================================")
    print(f"Verifying: {path}")
    print(f"Total Slides: {len(prs.slides)}")
    for idx, slide in enumerate(prs.slides):
        print(f"\n--- Slide {idx+1} ---")
        for shape in slide.shapes:
            if shape.has_text_frame:
                txt = " | ".join([p.text.strip() for p in shape.text_frame.paragraphs if p.text.strip()])
                if txt:
                    print(f"  [{shape.name}] ({shape.left.inches:.2f}, {shape.top.inches:.2f}, {shape.width.inches:.2f}, {shape.height.inches:.2f}): {txt[:90]}...")
