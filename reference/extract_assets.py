from pptx import Presentation
import os

prs = Presentation(r"c:\src\projects\SIH\voiceguard\reference\SIH2025-IDEA-Presentation-Format.pptx")
os.makedirs("reference/extracted_assets", exist_ok=True)

img_count = 0
for i, slide in enumerate(prs.slides):
    for shape in slide.shapes:
        if shape.shape_type == 13: # PICTURE
            img_count += 1
            image = shape.image
            image_bytes = image.blob
            ext = image.ext
            filename = f"reference/extracted_assets/slide_{i+1}_img_{img_count}.{ext}"
            with open(filename, "wb") as f:
                f.write(image_bytes)
            print(f"Extracted: {filename}, size: {len(image_bytes)} bytes, dims: {shape.width.inches:.2f}x{shape.height.inches:.2f}")

# Also check slide layouts and background
for i, slide in enumerate(prs.slides):
    bg = slide.background
    fill = bg.fill if bg else None
    print(f"Slide {i+1} background type: {fill.type if fill else 'None'}")
